import os
import glob
from pptx import Presentation
import subprocess

def get_video_thumb(filename, imgname):
    command = [
            "ffmpeg",
            '-y',  # overwrite output file if it exists
            '-loglevel', 'error',
            '-i', filename,
            '-vframes', '1',  # take only one frame
            imgname,
        ]
    subprocess.run(command, stdout=subprocess.PIPE)
    return imgname

def get_name(filename):
    pre, ext = os.path.splitext(filename)
    return pre.split(os.sep)[-1]

DIR = 'newMedia/videos/polarCurvesAnimation/1440p60/partial_movie_files'
scenes =  os.listdir(DIR)
SLD_BLANK = 6
prs = Presentation('templates/template.pptx')
slide_layout = prs.slide_layouts[SLD_BLANK]

selection = input("Scene Name: ")
for scene in scenes:
    if scene != selection:
        continue
    SCENE_DIR = os.path.join(DIR, scene)
    scene_clips = list(filter(lambda x: x.endswith('.mp4'), os.listdir(SCENE_DIR)))
    for file in sorted(scene_clips):
        file = os.path.join(SCENE_DIR, file)
        prs_ex = Presentation('templates/powerpoint.pptx')
        timing_ex = prs_ex.slides[0].element[2]
        slide = prs.slides.add_slide(slide_layout)
        
        thumb_file = os.path.join(SCENE_DIR, get_name(file) + ".png")
        get_video_thumb(file, thumb_file)
        
        clip = slide.shapes.add_movie(file, 0, 0, prs.slide_width, prs.slide_height, mime_type='video/mp4', poster_frame_image=thumb_file)
        id = clip.element[0][0].attrib.get("id")
        timing = timing_ex
        timing[0][0][0][0][0][0][0][0][0][1][0][0][1][0][0][1][0][0][1][0].attrib["spid"] = id
        timing[0][0][0][0][1][0][1][0].attrib["spid"] = id
        timing[0][0][0][0][2][0][0][0][0][0].attrib["spid"] = id
        timing[0][0][0][0][2][0][2][0][0][1][0][0][1][0][0][1][0][0][1][0].attrib["spid"] = id
        timing[0][0][0][0][2][1][0][0][0].attrib["spid"] = id
        slide.element[2] = timing
        prs.save(f"{SCENE_DIR}/{scene}.pptx")
    print(scene)
    